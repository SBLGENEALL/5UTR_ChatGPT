from __future__ import annotations

from pathlib import Path
import argparse
import gzip
import json
import shutil
import subprocess
import tempfile

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt


BASE = Path.cwd()
DEFAULT_SELECTED = BASE / "07_library_design/tables/selected_2000_50_100bp_cluster_diverse_evidence_balanced_library.csv"
DEFAULT_BASELINE = BASE / "07_library_design/tables/v1.3_selected_2000_library.csv"
DEFAULT_GENOME = BASE / "00_raw_data/01_ncbi_genome_annotation/GCF_003668045.1_CriGri-PICR_genomic.fna.gz"
DEFAULT_OUTDIR = BASE / "08_reports/v1.4_qc_audit"
SEQ_CANDIDATES = ["utr5_sequence_tss_corrected", "utr5_sequence", "sequence"]


def clean_seq(value):
    return "".join(c for c in str(value).upper().replace("U", "T") if c in "ACGTN")


def first_column(df, names, required=False):
    for name in names:
        if name in df.columns:
            return name
    if required:
        raise KeyError(f"Missing required column. Expected one of {names}; available={list(df.columns)}")
    return None


def numeric(df, column):
    if column and column in df.columns:
        return pd.to_numeric(df[column], errors="coerce")
    return pd.Series(np.nan, index=df.index)


def read_selected(path):
    df = pd.read_csv(path)
    seq_col = first_column(df, SEQ_CANDIDATES, required=True)
    df = df.copy()
    df["_sequence"] = df[seq_col].map(clean_seq)
    if "library_index" not in df.columns:
        df["library_index"] = np.arange(1, len(df) + 1)
    df["_query_id"] = "lib" + df["library_index"].astype(str)
    df["_length"] = numeric(df, first_column(df, ["length", "utr5_length_final", "utr5_length_tss_corrected"]))
    df["_length"] = df["_length"].fillna(df["_sequence"].str.len())
    df["_gc"] = numeric(df, first_column(df, ["gc_content", "gc"]))
    missing_gc = df["_gc"].isna()
    computed_gc = df["_sequence"].map(
        lambda s: (s.count("G") + s.count("C")) / len(s) if s else np.nan
    )
    df["_gc"] = df["_gc"].where(~missing_gc, computed_gc).astype(float)
    return df


def write_query_fasta(df, path):
    with path.open("w", encoding="ascii") as handle:
        for _, row in df.iterrows():
            handle.write(f">{row['_query_id']}\n{row['_sequence']}\n")


def parse_paf(path):
    columns = [
        "query_id", "query_length", "query_start", "query_end", "strand",
        "target_id", "target_length", "target_start", "target_end",
        "matches", "block_length", "mapq",
    ]
    rows = []
    with path.open("r", encoding="utf-8", errors="replace") as handle:
        for line in handle:
            fields = line.rstrip("\n").split("\t")
            if len(fields) < 12:
                continue
            row = dict(zip(columns, fields[:12]))
            for key in [
                "query_length", "query_start", "query_end", "target_length",
                "target_start", "target_end", "matches", "block_length", "mapq",
            ]:
                row[key] = int(row[key])
            row["identity"] = row["matches"] / row["block_length"] if row["block_length"] else np.nan
            row["coverage"] = (row["query_end"] - row["query_start"]) / row["query_length"] if row["query_length"] else np.nan
            row["bitscore"] = np.nan
            rows.append(row)
    return pd.DataFrame(rows)


def parse_blast(path):
    columns = [
        "query_id", "target_id", "pident", "alignment_length", "query_length",
        "mismatch", "gapopen", "query_start", "query_end", "target_start",
        "target_end", "evalue", "bitscore",
    ]
    if not path.exists() or path.stat().st_size == 0:
        return pd.DataFrame(columns=columns + ["identity", "coverage"])
    hits = pd.read_csv(path, sep="\t", names=columns)
    hits["identity"] = pd.to_numeric(hits["pident"], errors="coerce") / 100
    hits["coverage"] = (
        pd.to_numeric(hits["alignment_length"], errors="coerce")
        / pd.to_numeric(hits["query_length"], errors="coerce")
    )
    return hits


def materialize_genome(genome_path, workdir):
    if genome_path.suffix.lower() != ".gz":
        return genome_path
    target = workdir / genome_path.stem
    with gzip.open(genome_path, "rb") as source, target.open("wb") as dest:
        shutil.copyfileobj(source, dest)
    return target


def run_mapper(selected, genome_path, outdir, mapper):
    query_fasta = outdir / "selected_2000_queries.fasta"
    write_query_fasta(selected, query_fasta)
    mapping_path = outdir / ("cho_genome_mapping.paf" if mapper == "minimap2" else "cho_genome_mapping.tsv")

    if mapper == "minimap2":
        exe = shutil.which("minimap2")
        if not exe:
            raise RuntimeError("minimap2 was requested but is not available on PATH")
        with mapping_path.open("w", encoding="utf-8") as output:
            subprocess.run(
                [exe, "-x", "sr", "-N", "20", "--secondary=yes", str(genome_path), str(query_fasta)],
                check=True,
                stdout=output,
            )
        return parse_paf(mapping_path), mapping_path

    if mapper == "blastn":
        blastn = shutil.which("blastn")
        makeblastdb = shutil.which("makeblastdb")
        if not blastn or not makeblastdb:
            raise RuntimeError("blastn/makeblastdb were requested but are not available on PATH")
        with tempfile.TemporaryDirectory(prefix="v14_blast_", dir=outdir) as tmp:
            tmpdir = Path(tmp)
            genome_fasta = materialize_genome(genome_path, tmpdir)
            db = tmpdir / "cho_genome"
            subprocess.run([makeblastdb, "-in", str(genome_fasta), "-dbtype", "nucl", "-out", str(db)], check=True)
            outfmt = "6 qseqid sseqid pident length qlen mismatch gapopen qstart qend sstart send evalue bitscore"
            subprocess.run(
                [
                    blastn, "-task", "blastn-short", "-query", str(query_fasta), "-db", str(db),
                    "-outfmt", outfmt, "-max_target_seqs", "20", "-out", str(mapping_path),
                ],
                check=True,
            )
        return parse_blast(mapping_path), mapping_path

    raise ValueError(f"Unsupported mapper: {mapper}")


def choose_mapper(requested):
    if requested != "auto":
        return requested
    if shutil.which("minimap2"):
        return "minimap2"
    if shutil.which("blastn") and shutil.which("makeblastdb"):
        return "blastn"
    raise RuntimeError("Neither minimap2 nor blastn/makeblastdb is available on PATH")


def summarize_mapping(selected, hits, identity_min, coverage_min):
    if len(hits):
        hits = hits.copy()
        target_start = pd.to_numeric(hits.get("target_start"), errors="coerce").fillna(-1).astype(int)
        target_end = pd.to_numeric(hits.get("target_end"), errors="coerce").fillna(-1).astype(int)
        hits["_locus_key"] = (
            hits["target_id"].astype(str) + ":" + target_start.astype(str) + "-" + target_end.astype(str)
        )
        hits["qualifying_hit"] = (
            pd.to_numeric(hits["identity"], errors="coerce").ge(identity_min)
            & pd.to_numeric(hits["coverage"], errors="coerce").ge(coverage_min)
        )
        hits = hits.sort_values(
            ["query_id", "qualifying_hit", "identity", "coverage", "bitscore"],
            ascending=[True, False, False, False, False],
            na_position="last",
        )
    rows = []
    grouped = {key: group for key, group in hits.groupby("query_id")} if len(hits) else {}
    for _, row in selected.iterrows():
        query_id = row["_query_id"]
        group = grouped.get(query_id, pd.DataFrame())
        qualifying = group[group["qualifying_hit"]] if len(group) else group
        best = group.iloc[0] if len(group) else None
        qualifying_loci = qualifying["_locus_key"].nunique() if len(qualifying) else 0
        mapped = qualifying_loci > 0
        rows.append({
            "query_id": query_id,
            "library_index": row["library_index"],
            "gene_name": row.get("gene_name", ""),
            "library_group": row.get("library_group", ""),
            "mapped": mapped,
            "best_target": best["target_id"] if best is not None else "",
            "best_identity": float(best["identity"]) if best is not None else np.nan,
            "best_coverage": float(best["coverage"]) if best is not None else np.nan,
            "qualifying_locus_count": int(qualifying_loci),
            "unique_mapping": qualifying_loci == 1,
            "multi_mapping": qualifying_loci > 1,
            "suspected_non_CHO_or_hallucinated": not mapped,
        })
    return pd.DataFrame(rows)


def aggregate_qc(selected):
    data = selected.copy()
    data["length_bin"] = pd.cut(
        data["_length"], [49, 60, 70, 80, 90, 100],
        labels=["50-60", "61-70", "71-80", "81-90", "91-100"],
        include_lowest=True,
    )
    data["gc_bin"] = pd.cut(
        data["_gc"], [0.299, 0.40, 0.50, 0.60, 0.75],
        labels=["30-40%", "40-50%", "50-60%", "60-75%"],
        include_lowest=True,
    )
    metrics = {
        "robust_public_te_rank": first_column(data, ["robust_public_te_rank"]),
        "heavy_ensemble_score": first_column(data, ["heavy_ensemble_score", "automl_ensemble_score"]),
        "protein_abundance_rank": first_column(data, ["protein_abundance_rank"]),
        "protein_residual_rank": first_column(data, ["protein_residual_rank"]),
    }
    frames = {}
    for dimension in ["length_bin", "gc_bin"]:
        summary = pd.DataFrame({"bin": data[dimension].cat.categories.astype(str)})
        counts = data.groupby(dimension, observed=False).size().reindex(data[dimension].cat.categories)
        summary["selected_n"] = counts.values
        for label, column in metrics.items():
            values = numeric(data, column)
            means = values.groupby(data[dimension], observed=False).mean().reindex(data[dimension].cat.categories)
            summary[f"mean_{label}"] = means.values
        frames[dimension] = summary
    return data, frames


def save_charts(data, mapping_summary, outdir):
    chart_dir = outdir / "charts"
    chart_dir.mkdir(parents=True, exist_ok=True)
    chart_paths = []
    plot_specs = [
        ("_length", "robust_public_te_rank", "UTR length vs robust public TE rank", "length_vs_public_te.png"),
        ("_length", first_column(data, ["heavy_ensemble_score", "automl_ensemble_score"]), "UTR length vs heavy ensemble score", "length_vs_heavy_score.png"),
        ("_length", "protein_abundance_rank", "UTR length vs protein abundance rank", "length_vs_protein_abundance.png"),
        ("_length", "protein_residual_rank", "UTR length vs protein residual rank", "length_vs_protein_residual.png"),
        ("_gc", "robust_public_te_rank", "GC fraction vs robust public TE rank", "gc_vs_public_te.png"),
        ("_gc", "protein_abundance_rank", "GC fraction vs protein abundance rank", "gc_vs_protein_abundance.png"),
        ("_gc", "protein_residual_rank", "GC fraction vs protein residual rank", "gc_vs_protein_residual.png"),
    ]
    for xcol, ycol, title, filename in plot_specs:
        if not ycol or ycol not in data.columns:
            continue
        x = pd.to_numeric(data[xcol], errors="coerce")
        y = pd.to_numeric(data[ycol], errors="coerce")
        valid = x.notna() & y.notna()
        if not valid.any():
            continue
        fig, ax = plt.subplots(figsize=(7.4, 4.5))
        ax.scatter(x[valid], y[valid], s=14, alpha=0.45, color="#276FBF", edgecolors="none")
        ax.set_title(title)
        ax.set_xlabel("UTR length (nt)" if xcol == "_length" else "GC fraction")
        ax.set_ylabel(ycol.replace("_", " "))
        ax.grid(alpha=0.2)
        fig.tight_layout()
        path = chart_dir / filename
        fig.savefig(path, dpi=180)
        plt.close(fig)
        chart_paths.append(path)

    if len(mapping_summary):
        counts = pd.Series({
            "Unique": int(mapping_summary["unique_mapping"].sum()),
            "Multi": int(mapping_summary["multi_mapping"].sum()),
            "Unmapped": int((~mapping_summary["mapped"]).sum()),
        })
        fig, ax = plt.subplots(figsize=(6.5, 4.3))
        counts.plot(kind="bar", ax=ax, color=["#2A9D8F", "#E9C46A", "#D1495B"])
        ax.set_title("CHO genome mapping status")
        ax.set_ylabel("Selected UTR count")
        ax.tick_params(axis="x", rotation=0)
        ax.grid(axis="y", alpha=0.2)
        fig.tight_layout()
        path = chart_dir / "cho_genome_mapping_status.png"
        fig.savefig(path, dpi=180)
        plt.close(fig)
        chart_paths.append(path)
    return chart_paths


def compare_libraries(current, baseline_path):
    if not baseline_path or not baseline_path.exists():
        return pd.DataFrame([{
            "metric": "baseline_status",
            "v1.3": "not supplied",
            "v1.4": len(current),
            "delta": np.nan,
        }]), pd.DataFrame()
    baseline = read_selected(baseline_path)
    current_seq = set(current["_sequence"])
    baseline_seq = set(baseline["_sequence"])
    summary = pd.DataFrame([
        {"metric": "selected_n", "v1.3": len(baseline), "v1.4": len(current), "delta": len(current) - len(baseline)},
        {"metric": "shared_sequences", "v1.3": len(baseline_seq), "v1.4": len(current_seq & baseline_seq), "delta": np.nan},
        {"metric": "removed_from_v1.3", "v1.3": len(baseline_seq - current_seq), "v1.4": 0, "delta": np.nan},
        {"metric": "added_in_v1.4", "v1.3": 0, "v1.4": len(current_seq - baseline_seq), "delta": np.nan},
    ])
    detail = pd.DataFrame({
        "sequence": sorted(current_seq | baseline_seq),
    })
    detail["in_v1.3"] = detail["sequence"].isin(baseline_seq)
    detail["in_v1.4"] = detail["sequence"].isin(current_seq)
    detail["status"] = np.select(
        [detail["in_v1.3"] & detail["in_v1.4"], detail["in_v1.3"], detail["in_v1.4"]],
        ["shared", "removed", "added"],
        default="unknown",
    )
    return summary, detail


def read_optional_csv(path):
    return pd.read_csv(path) if path.exists() else pd.DataFrame()


def build_excel(out_path, summary, selected, length_summary, gc_summary, mapping, hits, comparison, comparison_detail, accession_audit, gene_audit):
    try:
        import openpyxl  # noqa: F401
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to create the v1.4 QC audit Excel workbook") from exc
    with pd.ExcelWriter(out_path, engine="openpyxl") as writer:
        summary.to_excel(writer, sheet_name="Executive Summary", index=False)
        length_summary.to_excel(writer, sheet_name="Length QC", index=False)
        gc_summary.to_excel(writer, sheet_name="GC QC", index=False)
        mapping.to_excel(writer, sheet_name="CHO Mapping Summary", index=False)
        hits.to_excel(writer, sheet_name="CHO Mapping Hits", index=False)
        comparison.to_excel(writer, sheet_name="v1.3 vs v1.4", index=False)
        comparison_detail.to_excel(writer, sheet_name="Sequence Comparison", index=False)
        accession_audit.to_excel(writer, sheet_name="Accession Gene Audit", index=False)
        gene_audit.to_excel(writer, sheet_name="Gene UTR Audit", index=False)
        selected.to_excel(writer, sheet_name="Selected Library", index=False)
        workbook = writer.book
        for sheet in workbook.worksheets:
            sheet.freeze_panes = "A2"
            sheet.auto_filter.ref = sheet.dimensions
            sheet.sheet_view.showGridLines = False
            for cell in sheet[1]:
                cell.font = openpyxl.styles.Font(bold=True, color="FFFFFF")
                cell.fill = openpyxl.styles.PatternFill("solid", fgColor="1F4E78")
            for column_cells in sheet.columns:
                width = min(45, max(10, max(len(str(cell.value or "")) for cell in column_cells[:200]) + 2))
                sheet.column_dimensions[column_cells[0].column_letter].width = width
        summary_sheet = workbook["Executive Summary"]
        summary_sheet["D1"] = "Mapping status"
        summary_sheet["E1"] = "Count"
        mapping_chart_rows = [
            ("Mapped", int(mapping["mapped"].sum())),
            ("Unmapped", int((~mapping["mapped"]).sum())),
            ("Unique", int(mapping["unique_mapping"].sum())),
            ("Multi", int(mapping["multi_mapping"].sum())),
        ]
        for row_index, (label, value) in enumerate(mapping_chart_rows, start=2):
            summary_sheet.cell(row=row_index, column=4, value=label)
            summary_sheet.cell(row=row_index, column=5, value=value)
        for cell in summary_sheet[1][3:5]:
            cell.font = openpyxl.styles.Font(bold=True, color="FFFFFF")
            cell.fill = openpyxl.styles.PatternFill("solid", fgColor="1F4E78")
        from openpyxl.chart import BarChart, Reference
        chart = BarChart()
        chart.type = "col"
        chart.style = 10
        chart.title = "CHO genome mapping status"
        chart.y_axis.title = "Selected UTR count"
        chart.legend = None
        chart.add_data(Reference(summary_sheet, min_col=5, min_row=1, max_row=5), titles_from_data=True)
        chart.set_categories(Reference(summary_sheet, min_col=4, min_row=2, max_row=5))
        chart.height = 7
        chart.width = 11
        summary_sheet.add_chart(chart, "D7")
        summary_sheet.column_dimensions["D"].width = 20
        summary_sheet.column_dimensions["E"].width = 12


def add_ppt_table(slide, rows, left, top, width, height):
    table = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height).table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 78, 120)
    return table


def build_ppt(out_path, summary, length_summary, mapping_summary, comparison, chart_paths):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to create the v1.4 QC summary presentation") from exc
    globals()["RGBColor"] = RGBColor
    globals()["Pt"] = Pt
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.8))
    title.text_frame.text = "5UTR Engineering v1.4 PR1 QC audit"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    rows = [["Metric", "Value"]] + [[r["metric"], r["value"]] for _, r in summary.head(10).iterrows()]
    add_ppt_table(slide, rows, Inches(0.8), Inches(1.6), Inches(5.7), Inches(4.9))
    if chart_paths:
        slide.shapes.add_picture(str(chart_paths[-1]), Inches(7.0), Inches(1.65), width=Inches(5.5))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.6))
    title.text_frame.text = "Length-wise evidence QC"
    title.text_frame.paragraphs[0].font.size = Pt(25)
    rows = [list(length_summary.columns)] + length_summary.fillna("").round(4).astype(str).values.tolist()
    add_ppt_table(slide, rows, Inches(0.7), Inches(1.25), Inches(6.0), Inches(4.8))
    length_chart = next((p for p in chart_paths if p.name == "length_vs_public_te.png"), None)
    if length_chart:
        slide.shapes.add_picture(str(length_chart), Inches(7.0), Inches(1.3), width=Inches(5.7))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.6))
    title.text_frame.text = "CHO genome-origin validation"
    title.text_frame.paragraphs[0].font.size = Pt(25)
    map_rows = [
        ["Status", "Count"],
        ["Mapped", int(mapping_summary["mapped"].sum())],
        ["Unmapped", int((~mapping_summary["mapped"]).sum())],
        ["Unique", int(mapping_summary["unique_mapping"].sum())],
        ["Multi", int(mapping_summary["multi_mapping"].sum())],
    ]
    add_ppt_table(slide, map_rows, Inches(0.8), Inches(1.4), Inches(4.0), Inches(3.0))
    map_chart = next((p for p in chart_paths if p.name == "cho_genome_mapping_status.png"), None)
    if map_chart:
        slide.shapes.add_picture(str(map_chart), Inches(5.4), Inches(1.15), width=Inches(6.8))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.6))
    title.text_frame.text = "v1.3 vs v1.4 selected library"
    title.text_frame.paragraphs[0].font.size = Pt(25)
    rows = [list(comparison.columns)] + comparison.fillna("").astype(str).values.tolist()
    add_ppt_table(slide, rows, Inches(0.8), Inches(1.4), Inches(8.5), Inches(3.8))
    prs.save(out_path)


def main():
    parser = argparse.ArgumentParser(description="Generate v1.4 PR1 QC audit outputs")
    parser.add_argument("--selected", type=Path, default=DEFAULT_SELECTED)
    parser.add_argument("--baseline-v1.3", dest="baseline_v13", type=Path, default=DEFAULT_BASELINE)
    parser.add_argument("--genome", type=Path, default=DEFAULT_GENOME)
    parser.add_argument("--mapper", choices=["auto", "minimap2", "blastn"], default="auto")
    parser.add_argument("--mapping-table", type=Path, default=None, help="existing PAF or BLAST outfmt 6 table")
    parser.add_argument("--outdir", type=Path, default=DEFAULT_OUTDIR)
    parser.add_argument("--identity-min", type=float, default=0.95)
    parser.add_argument("--coverage-min", type=float, default=0.90)
    parser.add_argument("--skip-ppt", action="store_true", help="skip PPT generation for dependency-limited validation runs")
    args = parser.parse_args()

    if not args.selected.exists():
        raise SystemExit(f"Missing selected library: {args.selected}")
    if not args.mapping_table and not args.genome.exists():
        raise SystemExit(f"Missing CHO genome: {args.genome}")
    args.outdir.mkdir(parents=True, exist_ok=True)

    selected = read_selected(args.selected)
    if args.mapping_table:
        hits = parse_paf(args.mapping_table) if args.mapping_table.suffix.lower() == ".paf" else parse_blast(args.mapping_table)
        mapper = "precomputed"
        mapping_path = args.mapping_table
    else:
        mapper = choose_mapper(args.mapper)
        hits, mapping_path = run_mapper(selected, args.genome, args.outdir, mapper)
    mapping = summarize_mapping(selected, hits, args.identity_min, args.coverage_min)
    data, aggregate = aggregate_qc(selected)
    comparison, comparison_detail = compare_libraries(selected, args.baseline_v13)

    accession_audit = read_optional_csv(
        BASE / "00_raw_data/05_cho_proteomics/Heffner_accession_to_gene_audit.csv"
    )
    gene_audit = read_optional_csv(
        BASE / "04_te_labeling/tables/proteomics_gene_to_corrected_utr_audit.csv"
    )
    chart_paths = save_charts(data, mapping, args.outdir)

    mean_identity = mapping["best_identity"].mean()
    mean_coverage = mapping["best_coverage"].mean()
    summary_values = {
        "selected_n": len(selected),
        "mapper": mapper,
        "mapping_table": mapping_path.name,
        "mapped_n": int(mapping["mapped"].sum()),
        "unmapped_n": int((~mapping["mapped"]).sum()),
        "unique_mapping_n": int(mapping["unique_mapping"].sum()),
        "multi_mapping_n": int(mapping["multi_mapping"].sum()),
        "suspected_non_CHO_or_hallucinated_n": int(mapping["suspected_non_CHO_or_hallucinated"].sum()),
        "mean_best_identity": float(mean_identity) if pd.notna(mean_identity) else None,
        "mean_best_coverage": float(mean_coverage) if pd.notna(mean_coverage) else None,
        "accession_audit_rows": len(accession_audit),
        "gene_UTR_audit_rows": len(gene_audit),
    }
    summary = pd.DataFrame([{"metric": key, "value": value} for key, value in summary_values.items()])

    hits.to_csv(args.outdir / "cho_genome_mapping_hits.csv", index=False)
    mapping.to_csv(args.outdir / "cho_genome_mapping_summary.csv", index=False)
    mapping[mapping["suspected_non_CHO_or_hallucinated"]].to_csv(
        args.outdir / "suspected_non_CHO_or_hallucinated_sequences.csv", index=False
    )
    aggregate["length_bin"].to_csv(args.outdir / "length_bin_qc.csv", index=False)
    aggregate["gc_bin"].to_csv(args.outdir / "gc_bin_qc.csv", index=False)
    comparison.to_csv(args.outdir / "v1.3_vs_v1.4_summary.csv", index=False)
    comparison_detail.to_csv(args.outdir / "v1.3_vs_v1.4_sequence_comparison.csv", index=False)
    (args.outdir / "v1.4_qc_audit_summary.json").write_text(
        json.dumps(summary_values, indent=2, allow_nan=False), encoding="utf-8"
    )

    build_excel(
        args.outdir / "v1.4_qc_audit.xlsx",
        summary, selected, aggregate["length_bin"], aggregate["gc_bin"],
        mapping, hits, comparison, comparison_detail, accession_audit, gene_audit,
    )
    if not args.skip_ppt:
        build_ppt(
            args.outdir / "v1.4_qc_audit_summary.pptx",
            summary, aggregate["length_bin"], mapping, comparison, chart_paths,
        )
    print("[SAVED]", args.outdir / "v1.4_qc_audit.xlsx")
    if not args.skip_ppt:
        print("[SAVED]", args.outdir / "v1.4_qc_audit_summary.pptx")
    print("[SAVED]", args.outdir / "charts")


if __name__ == "__main__":
    main()
